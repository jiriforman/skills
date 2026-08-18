#!/usr/bin/env python3
"""Render a .pptx with LibreOffice and crop regions for reuse in an HTML deck.

Rendering beats reading shape geometry: python-pptx reports GROUP children in the group's
own coordinate space, so positions inside a group are not slide coordinates.

    # 1. see what is on each slide, and dump every embedded image
    python3 crop_pptx.py inspect deck.pptx --out work/
    # 2. render every slide to PNG at 2.2x (≈3483x1960 for a 16:9 deck)
    python3 crop_pptx.py render deck.pptx --out work/
    # 3. crop a region as a fraction of the slide (left top right bottom)
    python3 crop_pptx.py crop work/slide_2.png 0.045 0.125 0.955 0.815 art.jpg --maxw 2600
"""
import argparse, os, subprocess, sys


def inspect(path, out):
    from pptx import Presentation
    prs = Presentation(path)
    os.makedirs(os.path.join(out, 'imgs'), exist_ok=True)
    W, H = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides, 1):
        print(f'\n=== SLIDE {i} ===')
        imgs = []

        def walk(shapes, depth=0):
            for sh in shapes:
                pct = lambda v, t: round(100 * v / t, 1) if v is not None else None
                print('  ' * depth + f'- {str(sh.shape_type):16} {sh.name[:26]:28} '
                      f'x={pct(sh.left, W)}% y={pct(sh.top, H)}% w={pct(sh.width, W)}%')
                if sh.shape_type == 6:            # GROUP — children use group coordinates!
                    walk(sh.shapes, depth + 1)
                elif sh.shape_type == 13 and hasattr(sh, 'image'):
                    imgs.append(sh)
                if sh.has_text_frame and sh.text_frame.text.strip():
                    print('  ' * depth + f'    TEXT: {sh.text_frame.text.strip()[:160]!r}')
                if getattr(sh, 'has_table', False) and sh.has_table:
                    for ri, row in enumerate(sh.table.rows):
                        cells = [c.text.replace("\n", " / ")[:60] for c in row.cells]
                        print('  ' * depth + f'    r{ri}: ' + ' || '.join(cells))
        walk(slide.shapes)
        for j, sh in enumerate(imgs):
            name = f's{i}_{j:02d}.{sh.image.ext}'
            open(os.path.join(out, 'imgs', name), 'wb').write(sh.image.blob)
        if imgs:
            print(f'  -> {len(imgs)} image(s) written to {out}/imgs')


def render(path, out, zoom=2.2):
    import fitz
    os.makedirs(out, exist_ok=True)
    subprocess.run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', out, path],
                   check=True, capture_output=True, timeout=300)
    pdf = os.path.join(out, os.path.splitext(os.path.basename(path))[0] + '.pdf')
    doc = fitz.open(pdf)
    for i, page in enumerate(doc, 1):
        f = os.path.join(out, f'slide_{i}.png')
        page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(f)
        print(f, )


def crop(src, box, dst, maxw=None, quality=88):
    from PIL import Image
    im = Image.open(src).convert('RGB')
    W, H = im.size
    l, t, r, b = box
    c = im.crop((int(l * W), int(t * H), int(r * W), int(b * H)))
    if maxw and c.size[0] > maxw:
        c = c.resize((maxw, int(maxw * c.size[1] / c.size[0])), Image.LANCZOS)
    if dst.lower().endswith(('.jpg', '.jpeg')):
        c.save(dst, 'JPEG', quality=quality, optimize=True)
    else:
        c.save(dst, 'PNG', optimize=True)
    print(f'{dst}: {c.size[0]}x{c.size[1]}  {os.path.getsize(dst)//1024} KB')


if __name__ == '__main__':
    p = argparse.ArgumentParser()
    sub = p.add_subparsers(dest='cmd', required=True)
    a = sub.add_parser('inspect'); a.add_argument('pptx'); a.add_argument('--out', default='work')
    b = sub.add_parser('render');  b.add_argument('pptx'); b.add_argument('--out', default='work')
    b.add_argument('--zoom', type=float, default=2.2)
    c = sub.add_parser('crop');    c.add_argument('png'); c.add_argument('box', nargs=4, type=float)
    c.add_argument('dst'); c.add_argument('--maxw', type=int); c.add_argument('--quality', type=int, default=88)
    args = p.parse_args()
    if args.cmd == 'inspect': inspect(args.pptx, args.out)
    elif args.cmd == 'render': render(args.pptx, args.out, args.zoom)
    else: crop(args.png, args.box, args.dst, args.maxw, args.quality)
